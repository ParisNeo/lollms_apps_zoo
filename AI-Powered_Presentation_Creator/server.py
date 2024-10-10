from fastapi import FastAPI, HTTPException, UploadFile, File
from pydantic import BaseModel
from pathlib import Path
import json
import asyncio
import aiohttp
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from fastapi.middleware.cors import CORSMiddleware

def create_presentation(json_structure):
    prs = Presentation()
    
    # Set presentation properties
    prs.core_properties.title = json_structure["presentation"]["title"]
    prs.core_properties.author = json_structure["presentation"].get("author", "LoLLMs")
    
    for slide_data in json_structure["presentation"]["slides"]:
        if slide_data["layout"] == "title_slide":
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data["title"]
            subtitle.text = slide_data.get("subtitle", "")
        
        elif slide_data["layout"] == "content_with_image":
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            tf = content.text_frame
            for item in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
            
            # Add image
            if "image" in slide_data:
                image_data = slide_data["image"]
                if image_data["type"] == "request":
                    # Generate image using AI
                    image_path = asyncio.run(generate_ai_image(image_data["prompt"]))
                    left = Inches(5)
                    top = Inches(2)
                    width = Inches(4)
                    height = Inches(3)
                    slide.shapes.add_picture(str(image_path), left, top, width, height)
        
        elif slide_data["layout"] == "two_content":
            slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
            title = slide.shapes.title
            left_content = slide.placeholders[1]
            right_content = slide.placeholders[2]
            
            title.text = slide_data["title"]
            
            # Left content
            left_title = left_content.text_frame.add_paragraph()
            left_title.text = slide_data["left_content"]["title"]
            left_title.font.bold = True
            
            if "image" in slide_data["left_content"]:
                image_data = slide_data["left_content"]["image"]
                if image_data["type"] == "user_provided":
                    # Use a placeholder for user-provided images
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(3))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
                    shape.line.color.rgb = RGBColor(100, 100, 100)
                    text_box = shape.text_frame
                    text_box.text = f"User Image: {image_data['description']}"
            
            # Right content
            right_title = right_content.text_frame.add_paragraph()
            right_title.text = slide_data["right_content"]["title"]
            right_title.font.bold = True
            
            if "image" in slide_data["right_content"]:
                image_data = slide_data["right_content"]["image"]
                if image_data["type"] == "generate":
                    # Generate image using AI
                    image_path = asyncio.run(generate_ai_image(image_data["prompt"]))
                    left = Inches(5.5)
                    top = Inches(2)
                    width = Inches(4)
                    height = Inches(3)
                    slide.shapes.add_picture(str(image_path), left, top, width, height)
    
    # Save the presentation
    output_path = Path("generated_presentations") / f"{json_structure['presentation']['title']}.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    return output_path

app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:9600","http://localhost:9605"],  # Allow CORS for localhost:9600
    allow_credentials=True,
    allow_methods=["*"],  # Allow all methods
    allow_headers=["*"],  # Allow all headers
)


class PresentationRequest(BaseModel):
    json_structure: str

class ImageUpload(BaseModel):
    slide_id: int
    image_description: str

@app.post("/generate_presentation")
async def generate_presentation(request: PresentationRequest):
    # Use LoLLMs to generate JSON structure based on user prompt
    json_structure = await generate_json_with_lollms(request.json_structure)
    
    # Create presentation
    output_path = create_presentation(json_structure)
    
    return {"message": "Presentation created", "path": str(output_path)}

@app.post("/upload_image")
async def upload_image(slide_id: int, file: UploadFile = File(...)):
    # Save uploaded image
    file_path = Path(f"uploads/image_{slide_id}.png")
    with file_path.open("wb") as buffer:
        buffer.write(await file.read())
    
    return {"message": "Image uploaded", "path": str(file_path)}

@app.post("/generate_image")
async def generate_image(prompt: str):
    # Integrate with AI image generation service
    image_path = await generate_ai_image(prompt)
    return {"message": "Image generated", "path": str(image_path)}

async def generate_json_with_lollms(prompt: str):
    # Implement LoLLMs interaction here
    # This is a placeholder implementation
    json_structure = {
        "presentation": {
            "title": "AI Generated Presentation",
            "slides": [
                {
                    "layout": "title_slide",
                    "title": "Welcome",
                    "subtitle": "Generated by LoLLMs"
                }
                # Add more slides based on the prompt
            ]
        }
    }
    return json_structure

async def generate_ai_image(prompt: str):
    # Implement AI image generation here
    # This is a placeholder implementation
    async with aiohttp.ClientSession() as session:
        async with session.post("https://ai-image-generator.example.com/generate", json={"prompt": prompt}) as response:
            if response.status == 200:
                image_data = await response.read()
                output_path = Path(f"generated_images/{prompt[:20]}.png")
                with output_path.open("wb") as f:
                    f.write(image_data)
                return output_path
            else:
                raise HTTPException(status_code=500, detail="Failed to generate image")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="localhost", port=8000)

