import pipmaster as pm

# List of required packages
required_packages = [
    "bs4",
    "selenium",
    "python-docx",
    "python-pptx",
    "PyPDF2",
    "aiohttp",
    "tiktoken",
    "fastapi",
    "uvicorn",
    "requests"
]

# Check and install missing packages
for package in required_packages:
    if not pm.is_installed(package):
        print(f"Installing {package}...")
        pm.install(package)
    else:
        print(f"{package} is already installed")

# Now proceed with the imports after ensuring all packages are installed
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Dict, Optional
import asyncio
from bs4 import BeautifulSoup
import json
import os
from pathlib import Path
from urllib.parse import urljoin, urlparse
import tiktoken
import logging
from lollms_client import LollmsClient
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import requests
from io import BytesIO
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile


lc = LollmsClient()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Data Models remain the same
class ScrapeRequest(BaseModel):
    urls: List[str]
    depth: int = 1

class ProcessRequest(BaseModel):
    text: str
    chunk_size: int = 4096

class SubjectUpdateRequest(BaseModel):
    chunk: str
    current_subjects: List[str]

# Storage configuration
STORAGE_DIR = Path("data")
STORAGE_DIR.mkdir(exist_ok=True)


# Helper Functions
def save_state(filename: str, data: dict):
    with open(os.path.join(STORAGE_DIR, filename), 'w', encoding='utf-8') as f:
        json.dump(data, f)

def load_state(filename: str) -> dict:
    try:
        with open(os.path.join(STORAGE_DIR, filename), 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        return {}


# File type handlers
def extract_text_from_pdf(content: bytes) -> str:
    try:
        pdf_file = BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return ""

def extract_text_from_docx(content: bytes) -> str:
    try:
        doc_file = BytesIO(content)
        doc = Document(doc_file)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        return ""

def extract_text_from_pptx(content: bytes) -> str:
    try:
        ppt_file = BytesIO(content)
        ppt = Presentation(ppt_file)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        return ""

class WebScraper:
    def __init__(self):
        self.chrome_options = Options()
        self.chrome_options.add_argument("--headless")  # Run in headless mode
        self.chrome_options.add_argument("--no-sandbox")
        self.chrome_options.add_argument("--disable-dev-shm-usage")
        
    async def scrape_url(self, url: str) -> str:
        file_extension = Path(urlparse(url).path).suffix.lower()
        
        if file_extension in ['.pdf', '.docx', '.pptx']:
            return await self._handle_document(url, file_extension)
        else:
            return await self._handle_webpage(url)

    async def _handle_document(self, url: str, extension: str) -> str:
        try:
            response = requests.get(url)
            content = response.content
            
            if extension == '.pdf':
                return extract_text_from_pdf(content)
            elif extension == '.docx':
                return extract_text_from_docx(content)
            elif extension == '.pptx':
                return extract_text_from_pptx(content)
            
        except Exception as e:
            logger.error(f"Error handling document {url}: {str(e)}")
            return ""

    async def _handle_webpage(self, url: str) -> str:
        driver = None
        try:
            service = Service()
            driver = webdriver.Chrome(service=service, options=self.chrome_options)
            driver.get(url)
            
            # Wait for the page to load (adjust timeout as needed)
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.TAG_NAME, "body"))
            )
            
            # Get the page source after JavaScript execution
            html_content = driver.page_source
            
            # Extract text using BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            text_elements = []
            for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'span']):
                text_elements.append(tag.get_text().strip())
            
            return "\n".join(filter(None, text_elements))
            
        except Exception as e:
            logger.error(f"Error scraping webpage {url}: {str(e)}")
            return ""
        finally:
            if driver:
                driver.quit()

    async def get_links(self, url: str) -> List[str]:
        driver = None
        try:
            service = Service()
            driver = webdriver.Chrome(service=service, options=self.chrome_options)
            driver.get(url)
            
            # Wait for the page to load
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.TAG_NAME, "body"))
            )
            
            # Get all links
            links = []
            elements = driver.find_elements(By.TAG_NAME, "a")
            base_domain = urlparse(url).netloc
            
            for element in elements:
                href = element.get_attribute("href")
                if href:
                    full_url = urljoin(url, href)
                    if urlparse(full_url).netloc == base_domain:
                        links.append(full_url)
            
            return links
            
        except Exception as e:
            logger.error(f"Error getting links from {url}: {str(e)}")
            return []
        finally:
            if driver:
                driver.quit()
                
def tokenize_text(text: str) -> List[int]:
    encoding = tiktoken.get_encoding("cl100k_base")
    return encoding.encode(text)

def chunk_tokens(tokens: List[int], chunk_size: int) -> List[List[int]]:
    return [tokens[i:i + chunk_size] for i in range(0, len(tokens), chunk_size)]

# Modified scrape endpoint
@app.post("/scrape")
async def scrape_urls(request: ScrapeRequest):
    print("Started scraping...")
    visited_urls = set()
    all_text = []
    scraper = WebScraper()
    
    async def crawl(url: str, depth: int):
        if depth < 0 or url in visited_urls:
            return
        
        visited_urls.add(url)
        text = await scraper.scrape_url(url)
        if text:
            all_text.append(text)
            
            if depth > 0:
                links = await scraper.get_links(url)
                tasks = [crawl(link, depth - 1) for link in links]
                await asyncio.gather(*tasks)
    
    tasks = [crawl(url, request.depth) for url in request.urls]
    await asyncio.gather(*tasks)
    
    combined_text = "\n\n".join(all_text)
    with (STORAGE_DIR / "raw_text.json").open('w', encoding='utf-8') as f:
        json.dump({"text": combined_text}, f)
    
    return {
        "status": "success", 
        "message": f"Scraped {len(visited_urls)} pages", 
        "data": combined_text
    }

@app.post("/process")
async def process_text(request: ProcessRequest):
    tokens = tokenize_text(request.text)
    chunks = chunk_tokens(tokens, request.chunk_size)
    
    # Save chunks
    save_state("chunks.json", {"chunks": chunks})
    
    return {"status": "success", "chunks_count": len(chunks)}

@app.post("/update_subjects")
async def update_subjects(request: SubjectUpdateRequest):
    current_state = load_state("subjects.json")
    
    # Use Lollms to process the chunk and update subjects
    # This is a placeholder for the actual LLM integration
    # You would call your LLM here with the appropriate prompt
    
    # For now, we'll just store the chunks and subjects
    if "processed_chunks" not in current_state:
        current_state["processed_chunks"] = []
    current_state["processed_chunks"].append(request.chunk)
    current_state["current_subjects"] = request.current_subjects
    
    save_state("subjects.json", current_state)
    
    return {"status": "success", "current_subjects": request.current_subjects}

@app.get("/status")
async def get_status():
    try:
        raw_text = load_state("raw_text.json")
        chunks = load_state("chunks.json")
        subjects = load_state("subjects.json")
        
        return {
            "has_raw_text": bool(raw_text),
            "chunks_count": len(chunks.get("chunks", [])),
            "processed_chunks": len(subjects.get("processed_chunks", [])),
            "subjects": subjects.get("current_subjects", [])
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/rewrite_subject/{subject}")
async def rewrite_subject(subject: str):
    subjects_data = load_state("subjects.json")
    if subject not in subjects_data.get("current_subjects", []):
        raise HTTPException(status_code=404, detail="Subject not found")
    
    # Here you would use Lollms to rewrite the subject content
    # This is a placeholder for the actual LLM integration
    
    return {"status": "success", "subject": subject}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="localhost", port=8000)
